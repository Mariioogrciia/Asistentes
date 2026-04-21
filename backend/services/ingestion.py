"""Service: document ingestion pipeline.

Pipeline: upload file → extract text → chunk → embed → store in Supabase.
"""
from openai import AzureOpenAI

import uuid
from pathlib import Path

from loguru import logger
from supabase import Client

from backend.ai import get_embedding_client
from backend.config import get_settings

# Supported MIME types → file extensions
SUPPORTED_TYPES: dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
    "text/markdown": "md",
}


# ── Text extraction ────────────────────────────────────────────────────────────

def _extract_text(content: bytes, file_type: str, filename: str) -> str:
    """Extract plain text from a file based on its type."""
    if file_type == "pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    if file_type == "docx":
        import io
        from docx import Document
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if file_type == "pptx":
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)

    # txt / md → decode as UTF-8 (fallback latin-1)
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


# ── Chunking ───────────────────────────────────────────────────────────────────

def _split_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Split text into overlapping chunks using LangChain's splitter."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return [c for c in splitter.split_text(text) if c.strip()]


# ── Embedding in batches ────────────────────────────────────────────────────────

def _embed_chunks(
    chunks: list[str],
    batch_size: int = 16,
) -> list[list[float]]:
    """Embed chunks in batches using the embedding-specific Azure OpenAI client."""
    settings = get_settings()
    client = get_embedding_client()  # accessed at call-time, always initialized
    vectors: list[list[float]] = []
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i : i + batch_size]
        logger.debug("Embedding batch {}/{} ({} chunks)", i // batch_size + 1, -(-len(chunks) // batch_size), len(batch))
        response = client.embeddings.create(
            model=settings.azure_deployment_embedding,
            input=batch,
        )
        vectors.extend(item.embedding for item in response.data)
    return vectors


# ── Main ingestion function ────────────────────────────────────────────────────

def ingest_document(
    *,
    content: bytes,
    filename: str,
    content_type: str,
    assistant_id: uuid.UUID,
    assistant_name: str,
    db: Client,
    ai_client: AzureOpenAI,
) -> dict:
    """Full ingestion pipeline for a single document.

    Returns the created document record from Supabase.
    Raises ValueError for unsupported file types.
    """
    settings = get_settings()

    # ── 1. Validate file type ──────────────────────────────────────────────
    file_type = SUPPORTED_TYPES.get(content_type)
    if file_type is None:
        # Try guessing from filename extension
        ext = Path(filename).suffix.lstrip(".").lower()
        if ext in SUPPORTED_TYPES.values():
            file_type = ext
        else:
            raise ValueError(
                f"Unsupported file type: {content_type}. "
                f"Supported: {list(SUPPORTED_TYPES.keys())}"
            )

    doc_id = uuid.uuid4()
    logger.info(
        "Ingestion started doc_id={} filename={} type={} size={}B",
        doc_id,
        filename,
        file_type,
        len(content),
    )

    # ── 2. Create document record (status=processing) ──────────────────────
    # Sanitize assistant name and filename for Storage: only ASCII alphanumeric + . _ -
    # Supabase Storage rejects non-ASCII chars (e.g. ñ, á, é) in the key path
    clean_assistant_name = "".join(
        c if (c.isascii() and c.isalnum()) or c in "._-" else "_"
        for c in assistant_name
    )
    clean_filename = "".join(
        c if (c.isascii() and c.isalnum()) or c in "._-" else "_"
        for c in filename
    )
    storage_path = f"{clean_assistant_name}/{doc_id}/{clean_filename}"
    
    doc_record = {
        "id": str(doc_id),
        "assistant_id": str(assistant_id),
        "filename": filename, # Keep original filename in DB
        "file_type": file_type,
        "storage_path": storage_path,
        "size_bytes": len(content),
        "chunk_count": 0,
        "status": "processing",
    }
    db.table("documents").insert(doc_record).execute()

    try:
        # ── 3. Upload original file to Supabase Storage ────────────────────
        logger.debug("Uploading file to Storage path={}", storage_path)
        db.storage.from_(settings.supabase_bucket).upload(
            path=storage_path,
            file=content,
            file_options={"content-type": content_type},
        )

        # ── 4. Extract text ────────────────────────────────────────────────
        logger.debug("Extracting text from {}", filename)
        text = _extract_text(content, file_type, filename)
        if not text.strip():
            raise ValueError(f"No text could be extracted from {filename}")

        # ── 5. Chunk ───────────────────────────────────────────────────────
        chunks = _split_text(text, settings.chunk_size, settings.chunk_overlap)
        logger.info("doc_id={} — {} chunks produced", doc_id, len(chunks))

        # ── 6. Embed ───────────────────────────────────────────────────────
        vectors = _embed_chunks(chunks)

        # ── 7. Store chunks with vectors ───────────────────────────────────
        chunk_rows = [
            {
                "assistant_id": str(assistant_id),
                "document_id": str(doc_id),
                "content": chunk,
                "chunk_index": idx,
                # pgvector expects the vector as a plain Python list —
                # supabase-py serializes it correctly via postgrest-py
                "embedding": vector,
                "metadata": {"chunk_index": idx},
            }
            for idx, (chunk, vector) in enumerate(zip(chunks, vectors))
        ]
        # Insert in batches of 100 to avoid request size limits
        for i in range(0, len(chunk_rows), 100):
            db.table("chunks").insert(chunk_rows[i : i + 100]).execute()

        # ── 8. Mark document as ready ──────────────────────────────────────
        result = (
            db.table("documents")
            .update({"status": "ready", "chunk_count": len(chunks)})
            .eq("id", str(doc_id))
            .execute()
        )
        logger.info("Ingestion complete doc_id={} chunks={}", doc_id, len(chunks))
        return result.data[0]

    except Exception as exc:
        logger.error("Ingestion failed doc_id={} error={}", doc_id, exc)
        db.table("documents").update({"status": "error"}).eq("id", str(doc_id)).execute()
        raise
